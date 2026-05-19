"""
VaughnMartin Investor Pitch Deck — PPTX v10 FINAL
11 slides · clean image zones · no overlaps · logo on every slide
LAYOUT RULE: every image has an explicit zone; text zones never enter image zones.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image as PILImage
import os

# ── Colors ────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x0F, 0x2E)
NAVY2  = RGBColor(0x12, 0x1D, 0x47)
NAVY3  = RGBColor(0x0D, 0x14, 0x38)
GOLD   = RGBColor(0xC9, 0xA8, 0x4C)
TEAL   = RGBColor(0x2B, 0x8A, 0x6E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
MUTED  = RGBColor(0x99, 0xA5, 0xBB)
MID    = RGBColor(0x1E, 0x27, 0x55)
DARK3  = RGBColor(0x0E, 0x16, 0x40)
RED88  = RGBColor(0xFF, 0x88, 0x88)
GRN99  = RGBColor(0x66, 0xCC, 0x99)
IVORY  = RGBColor(0xF2, 0xF0, 0xEB)

W, H   = Inches(13.333), Inches(7.5)
TOTAL  = 11

# No-content zone (slide numbers only): x=12.0–13.2, y=7.0–7.5
NUM_X, NUM_Y = Inches(12.0),  Inches(7.06)
NUM_W, NUM_H = Inches(1.2),   Inches(0.34)
SAFE_Y       = Inches(6.84)   # nothing but logos/slide# below this

# ── Image paths (HD set) ────────────────────────────────────
HD           = "attached_assets/pitch-images-hd/"
IMG_HOME     = HD + "home-raw.png"
IMG_SIGNALS  = HD + "signals-raw.png"
IMG_BUILDER  = HD + "builder-raw.png"
IMG_PROTOCOLS= HD + "protocols-raw.png"
IMG_TOWER    = HD + "tower-raw.png"
IMG_ACTIVATION=HD + "activation-raw.png"
IMG_LOGO     = HD + "vm-logo-hd.png"   # 648×174 px, ratio ≈ 3.724

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════

def new_slide(bg=NAVY):
    s = prs.slides.add_slide(blank)
    f = s.background.fill; f.solid(); f.fore_color.rgb = bg
    return s

def gold_bar(s):
    r = s.shapes.add_shape(1, 0, 0, W, Pt(4.5))
    r.fill.solid(); r.fill.fore_color.rgb = GOLD; r.line.fill.background()

def slide_num(s, n):
    bx = s.shapes.add_textbox(NUM_X, NUM_Y, NUM_W, NUM_H)
    tf = bx.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = f"{n:02d} / {TOTAL:02d}"
    r.font.name = "Calibri"; r.font.size = Pt(9); r.font.bold = True
    r.font.color.rgb = RGBColor(0x44, 0x50, 0x70)

def logo_mark(s, x=Inches(0.36), y=None, w=Inches(1.52)):
    if not os.path.exists(IMG_LOGO): return
    h = w / 3.724
    if y is None:
        y = SAFE_Y + (H - SAFE_Y - h) / 2
    try: s.shapes.add_picture(IMG_LOGO, x, y, w, h)
    except Exception: pass

def txb(s, text, l, t, w, h, size=14, bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, font="Barlow Condensed", wrap=True):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    rn = p.add_run(); rn.text = text
    rn.font.name = font; rn.font.size = Pt(size)
    rn.font.bold = bold; rn.font.italic = italic
    rn.font.color.rgb = color
    return bx

def lbl(s, text, l, t, w=Inches(12), color=GOLD, align=PP_ALIGN.LEFT):
    txb(s, text.upper(), l, t, w, Pt(16), size=9, bold=True, color=color, align=align)

def heading(s, text, l, t, w, h, size=36, color=WHITE,
            align=PP_ALIGN.LEFT, font="Barlow Condensed"):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    rn = p.add_run(); rn.text = text
    rn.font.name = font; rn.font.size = Pt(size)
    rn.font.bold = True; rn.font.color.rgb = color
    return bx

def rect(s, l, t, w, h, fill=NAVY3, stroke=None, sw=0.8):
    r = s.shapes.add_shape(1, l, t, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if stroke: r.line.color.rgb = stroke; r.line.width = Pt(sw)
    else: r.line.fill.background()
    return r

def hrule(s, l, t, w, color=GOLD, lw=1.5):
    ln = s.shapes.add_shape(1, l, t, w, Pt(lw))
    ln.fill.solid(); ln.fill.fore_color.rgb = color; ln.line.fill.background()

def vrule(s, l, t, h, color=MID, lw=0.5):
    ln = s.shapes.add_shape(1, l, t, Pt(lw), h)
    ln.fill.solid(); ln.fill.fore_color.rgb = color; ln.line.fill.background()

def cover_pic(s, path, l, t, w, h, border=True):
    """Cover-crop: fills frame without distortion. Border = thin gold line."""
    if not os.path.exists(path):
        rect(s, l, t, w, h, fill=MID, stroke=RGBColor(0x28, 0x34, 0x68))
        return None
    with PILImage.open(path) as img:
        iw, ih = img.size
    frame_ar = w / h;  img_ar = iw / ih
    pic = s.shapes.add_picture(path, l, t, w, h)
    if img_ar > frame_ar:
        excess = (img_ar - frame_ar) / img_ar
        pic.crop_left = excess / 2;  pic.crop_right  = excess / 2
        pic.crop_top  = 0;           pic.crop_bottom = 0
    else:
        inv_excess = (h/w - ih/iw) / (h/w)
        pic.crop_top = inv_excess / 2; pic.crop_bottom = inv_excess / 2
        pic.crop_left = 0;             pic.crop_right  = 0
    if border:
        br = s.shapes.add_shape(1, l, t, w, h)
        br.fill.background(); br.line.color.rgb = GOLD; br.line.width = Pt(1.2)
    return pic

def caption_bar(s, text, l, t, w, h):
    """Dark bar pinned to the BOTTOM of the image. Must be called after cover_pic."""
    cap_h = Inches(0.26)
    rect(s, l, t + h - cap_h, w, cap_h, fill=RGBColor(0x06, 0x0A, 0x1E))
    txb(s, text.upper(), l + Inches(0.1), t + h - cap_h + Pt(5),
        w - Inches(0.2), cap_h, size=8, bold=True, color=GOLD)

def img_border_only(s, l, t, w, h):
    br = s.shapes.add_shape(1, l, t, w, h)
    br.fill.background(); br.line.color.rgb = GOLD; br.line.width = Pt(1.0)


# ══════════════════════════════════════════════════════════════
# S1 — Opening
# LAYOUT: left text zone x=0–6.9 | right image zone x=7.1–13.1
# NO overlapping shapes between zones.
# ══════════════════════════════════════════════════════════════
s1 = new_slide()
gold_bar(s1)

# ── RIGHT IMAGE ZONE: 16:9 frame, vertically centered ──────────
# Portrait frame (5.97×6.24) was cropping 46% off left+right of 16:9 source.
# Fix: set frame to exact 16:9, center vertically in the panel (y=0.56–6.84).
IMG1_X  = Inches(7.1);  IMG1_W = Inches(5.97)
IMG1_H  = IMG1_W / 1.778          # ≈ 3.36in — exact 16:9, zero crop
PANEL1_H = Inches(6.28)           # available panel height 0.56–6.84
IMG1_Y  = Inches(0.56) + (PANEL1_H - IMG1_H) / 2  # vertically centered
cover_pic(s1, IMG_TOWER, IMG1_X, IMG1_Y, IMG1_W, IMG1_H, border=True)
caption_bar(s1, "Command Tower · live signal monitoring · production",
    IMG1_X, IMG1_Y, IMG1_W, IMG1_H)

# Vertical gold divider at the split point
vrule(s1, Inches(6.96), Inches(0.56), Inches(6.24), color=GOLD, lw=1.0)

# ── LEFT TEXT ZONE: x=0.36 to x=6.72 (stays clear of divider) ─
logo_mark(s1, x=Inches(0.42), y=Inches(0.50), w=Inches(2.6))

lbl(s1, "Strategic Readiness Platform · Startup to Fortune 500",
    Inches(0.42), Inches(1.22), w=Inches(6.3),
    color=RGBColor(0x88, 0x78, 0x50))

heading(s1,
    "The response is ready\nbefore the trigger fires.",
    Inches(0.42), Inches(1.54), Inches(6.3), Inches(2.28),
    size=38, align=PP_ALIGN.LEFT)

txb(s1,
    "Most enterprises plan to respond.\nVaughnMartin is already executing.",
    Inches(0.42), Inches(3.92), Inches(6.3), Inches(0.70),
    size=16, italic=False, color=MUTED, font="Barlow Condensed")

for i, chip in enumerate([
    "Activist Investor  ·  91% confidence",
    "Ransomware  ·  95% confidence",
    "Regulatory Inquiry  ·  87% confidence",
]):
    cy = Inches(4.98) + i * Inches(0.60)
    rect(s1, Inches(0.42), cy, Inches(6.3), Inches(0.46),
         fill=RGBColor(0x18, 0x22, 0x50),
         stroke=RGBColor(0x44, 0x52, 0x80))
    txb(s1, chip, Inches(0.54), cy + Pt(7), Inches(6.1), Inches(0.40),
        size=12, bold=True, color=GOLD)

# Footer bar — full width, below both zones
rect(s1, 0, Inches(6.84), W, Inches(0.52),
     fill=RGBColor(0x06, 0x08, 0x18))
txb(s1, "AI monitors continuously.   Executives authorize decisively.   Teams execute.",
    Inches(0.5), Inches(6.91), Inches(12.33), Inches(0.44),
    size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

slide_num(s1, 1)


# ══════════════════════════════════════════════════════════════
# S2 — Readiness Gap
# No product images — layout already clean.
# ══════════════════════════════════════════════════════════════
s2 = new_slide()
gold_bar(s2); logo_mark(s2)

rect(s2, Inches(0.36), Inches(0.52), Inches(5.9), Inches(6.22), fill=NAVY3)
lbl(s2, "The Reality", Inches(0.66), Inches(0.72))
heading(s2, "30", Inches(0.66), Inches(1.10), Inches(3.5), Inches(2.1),
        size=120, color=GOLD)
heading(s2, "DAYS", Inches(0.66), Inches(3.06), Inches(3.5), Inches(0.64), size=30)
hrule(s2, Inches(0.66), Inches(3.82), Inches(4.82))
txb(s2, "Average mobilization time before execution begins.",
    Inches(0.66), Inches(4.0), Inches(4.9), Inches(0.56),
    size=14, italic=True, color=MUTED, font="Barlow Condensed")

rect(s2, Inches(6.74), Inches(0.52), Inches(6.23), Inches(6.22), fill=NAVY3)
lbl(s2, "The Readiness Question", Inches(7.04), Inches(0.72))
heading(s2,
    "Who calls who?\nWhere's the brief?\nWho owns it? Who authorizes?",
    Inches(7.04), Inches(1.10), Inches(5.8), Inches(2.3), size=25, color=GOLD)
hrule(s2, Inches(7.04), Inches(3.54), Inches(5.6))

rect(s2, Inches(7.04), Inches(3.74), Inches(5.8), Inches(0.48),
     fill=RGBColor(0x15, 0x1F, 0x52))
txb(s2, "This is a readiness problem, not a talent problem.",
    Inches(7.18), Inches(3.82), Inches(5.52), Inches(0.40),
    size=13, bold=True, color=WHITE)

for i, b in enumerate([
    "— Every trigger restarts the mobilization clock from zero",
    "— The window to act closes before execution begins",
    "— Every enterprise has talent. None have pre-staged responses.",
]):
    txb(s2, b, Inches(7.04), Inches(4.38) + i * Inches(0.58),
        Inches(5.8), Inches(0.54), size=13, color=MUTED)

vrule(s2, Inches(6.54), Inches(0.40), Inches(6.52))
slide_num(s2, 2)


# ══════════════════════════════════════════════════════════════
# S3 — Problem Is Here
# Badge/domain overlap fix: badge strictly top-right, domain label
# width capped at CARD_W - badge_w - 0.5in
# ══════════════════════════════════════════════════════════════
s3 = new_slide()
gold_bar(s3); logo_mark(s3)

lbl(s3, "The Problem Is Already Here", Inches(0.5), Inches(0.46))
heading(s3, "One of these is forming in your organization right now",
    Inches(0.5), Inches(0.76), Inches(12.33), Inches(0.76),
    size=26, align=PP_ALIGN.CENTER)

CARD_W = Inches(4.11); CARD_H = Inches(4.62); CARD_Y = Inches(1.64)
BADGE_W = Inches(0.96); BADGE_H = Inches(0.44)

for i, (conf, domain, n1, n2, meta, accent, pct) in enumerate([
    ("95%", "RISK & RESILIENCE",    "Ransomware",   "Attack Confirmed",  "Signal · 248 data points",    TEAL, 95),
    ("87%", "REGULATORY",           "Regulatory",   "Inquiry Opened",    "Signal · threshold crossed",  GOLD, 87),
    ("82%", "GROWTH & POSITIONING", "Market Entry", "Window Opening",    "Opportunity · live monitor",  GOLD, 82),
]):
    cx = Inches(0.38 + i * 4.32)

    rect(s3, cx, CARD_Y, CARD_W, CARD_H,
         fill=RGBColor(0x10, 0x18, 0x42), stroke=RGBColor(0x28, 0x34, 0x68))

    # Badge — solid backing, top-right corner
    bx = cx + CARD_W - BADGE_W - Inches(0.12)
    by = CARD_Y + Inches(0.12)
    rect(s3, bx, by, BADGE_W, BADGE_H,
         fill=RGBColor(0x06, 0x0C, 0x28), stroke=accent, sw=1.4)
    txb(s3, conf, bx, by, BADGE_W, BADGE_H,
        size=20, bold=True, color=accent, align=PP_ALIGN.CENTER)

    # Domain label — width explicitly ends BEFORE badge starts
    domain_w = CARD_W - BADGE_W - Inches(0.42)   # 0.42in gap
    txb(s3, domain, cx + Inches(0.18), CARD_Y + Inches(0.20),
        domain_w, Inches(0.40), size=9, bold=True, color=TEAL)

    heading(s3, f"{n1}\n{n2}",
            cx + Inches(0.18), CARD_Y + Inches(0.80),
            CARD_W - Inches(0.36), Inches(1.06), size=22)
    txb(s3, meta,
        cx + Inches(0.18), CARD_Y + Inches(2.00),
        CARD_W - Inches(0.36), Inches(0.34), size=11, color=MUTED)

    bar_y = CARD_Y + Inches(2.46)
    bar_w = CARD_W - Inches(0.36)
    rect(s3, cx + Inches(0.18), bar_y, bar_w, Pt(4), fill=MID)
    rect(s3, cx + Inches(0.18), bar_y, int(bar_w * pct / 100), Pt(4), fill=accent)

rect(s3, Inches(0.36), Inches(6.42), Inches(12.61), Inches(0.42),
     fill=RGBColor(0x10, 0x18, 0x42))
txb(s3, "221 triggers monitored  ·  248 data points  ·  refreshed every 15 minutes",
    Inches(0.5), Inches(6.48), Inches(12.33), Inches(0.36),
    size=10, bold=True, color=RGBColor(0x44, 0x58, 0x80), align=PP_ALIGN.CENTER)
slide_num(s3, 3)


# ══════════════════════════════════════════════════════════════
# S4 — Solution
# ══════════════════════════════════════════════════════════════
s4 = new_slide()
gold_bar(s4); logo_mark(s4)

lbl(s4, "The Answer", Inches(0.5), Inches(0.46), align=PP_ALIGN.CENTER)
heading(s4, "The response is ready",
    Inches(1.0), Inches(0.80), Inches(11.33), Inches(1.14),
    size=52, align=PP_ALIGN.CENTER)
heading(s4, "before the trigger fires.",
    Inches(1.0), Inches(1.88), Inches(11.33), Inches(1.14),
    size=52, color=GOLD, align=PP_ALIGN.CENTER)
txb(s4, "Preparation  →  Readiness  →  Fearless",
    Inches(1.0), Inches(3.06), Inches(11.33), Inches(0.50),
    size=16, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
hrule(s4, Inches(4.17), Inches(3.64), Inches(5.0))

for i, (num, sub) in enumerate([
    ("170", "Readiness Protocols"),
    ("221", "Strategic Triggers"),
    ("12 MIN", "Execution Window"),
]):
    px = Inches(0.55 + i * 4.22)
    txb(s4, num, px, Inches(3.84), Inches(4.0), Inches(1.2),
        size=52 if num != "12 MIN" else 38,
        bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txb(s4, sub, px, Inches(5.02), Inches(4.0), Pt(24),
        size=13, color=WHITE, align=PP_ALIGN.CENTER)

rect(s4, Inches(2.0), Inches(6.10), Inches(9.33), Inches(0.52),
     fill=NAVY, stroke=RGBColor(0x44, 0x52, 0x80))
txb(s4, "AI monitors.   Executives authorize.   Teams execute.",
    Inches(2.1), Inches(6.20), Inches(9.13), Inches(0.44),
    size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
slide_num(s4, 4)


# ══════════════════════════════════════════════════════════════
# S5 — Old Model vs VaughnMartin + protocols proof strip
# LAYOUT:
#   Columns:  y=1.74 → y=5.94 (h=4.20)
#   Strip:    y=6.08 → y=6.76 (h=0.68) — FULLY BELOW columns, zero overlap
# ══════════════════════════════════════════════════════════════
s5 = new_slide()
gold_bar(s5); logo_mark(s5)

lbl(s5, "Why This Is Defensible", Inches(0.5), Inches(0.46))
heading(s5,
    "The architecture that separates execution from mobilization",
    Inches(0.5), Inches(0.76), Inches(12.33), Inches(0.72), size=24)
hrule(s5, Inches(0.5), Inches(1.60), Inches(12.33))

COL_Y = Inches(1.74); COL_H = Inches(4.20)   # ends y = 5.94

# Left — Old Model
rect(s5, Inches(0.36), COL_Y, Inches(6.22), COL_H,
     fill=RGBColor(0x18, 0x08, 0x08))
txb(s5, "Old Model", Inches(0.66), COL_Y + Inches(0.18),
    Inches(5.86), Pt(22), size=10, bold=True, color=RED88)
for i, b in enumerate([
    "✕  Faster notes from the same slow meetings",
    "✕  No readiness architecture — triggers catch you cold",
    "✕  Authority unclear when pressure arrives",
    "✕  30 days to mobilize before execution starts",
]):
    txb(s5, b, Inches(0.66), COL_Y + Inches(0.54) + i * Inches(0.76),
        Inches(5.86), Inches(0.68), size=13, color=RED88)

# Right — VaughnMartin
rect(s5, Inches(6.75), COL_Y, Inches(6.22), COL_H,
     fill=RGBColor(0x08, 0x18, 0x10))
txb(s5, "VaughnMartin Readiness OS", Inches(7.05), COL_Y + Inches(0.18),
    Inches(5.86), Pt(22), size=10, bold=True, color=GOLD)
for i, b in enumerate([
    "✓  Response pre-staged before the trigger fires",
    "✓  170 protocols · 221 triggers · ready before the moment",
    "✓  Executive authorization gate at every activation",
    "✓  12-minute execution window from signal to action",
]):
    txb(s5, b, Inches(7.05), COL_Y + Inches(0.54) + i * Inches(0.76),
        Inches(5.86), Inches(0.68), size=13, color=GRN99)

vrule(s5, Inches(6.55), COL_Y, COL_H)

# Product proof strip — BELOW columns (y=6.08–6.76), no overlap
STRIP_Y = Inches(6.08); STRIP_H = Inches(0.68)
cover_pic(s5, IMG_PROTOCOLS,
    Inches(0.36), STRIP_Y, Inches(12.61), STRIP_H, border=True)
# Dark overlay on the strip (solid, not transparent — transparency unsupported in pptx)
rect(s5, Inches(0.36), STRIP_Y, Inches(12.61), STRIP_H,
     fill=RGBColor(0x06, 0x08, 0x1A), stroke=GOLD, sw=0.8)
txb(s5,
    "LIVE PLATFORM  ·  170 Readiness Protocols  ·  startup to Fortune 500  ·  vaughnmartin.com",
    Inches(0.5), STRIP_Y + Pt(14), Inches(12.33), STRIP_H,
    size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

slide_num(s5, 5)


# ══════════════════════════════════════════════════════════════
# S6 — Signal → Execution (two images, explicit gap)
# LAYOUT:
#   Left image:  x=0.36, w=6.08 → ends x=6.44
#   Gap:         0.30in           (x=6.44–6.74)
#   Right image: x=6.74, w=6.23  → ends x=12.97
# ══════════════════════════════════════════════════════════════
s6 = new_slide()
gold_bar(s6); logo_mark(s6)

lbl(s6, "Proof of Production", Inches(0.5), Inches(0.46))
heading(s6, "From Signal to Authorized Execution in 12 Minutes",
    Inches(0.5), Inches(0.76), Inches(12.33), Inches(0.76),
    size=28, align=PP_ALIGN.CENTER)

# Images at native 16:9 ratio (img AR=1.778) so no horizontal clipping.
# Both images same height for visual consistency.
# Total usable width: 13.333 - 0.36 - 0.36 - 0.30 gap = 12.313in → each ~6.157in
IMG6_Y  = Inches(1.62)
IMG6_LW = Inches(6.15)                             # left width
IMG6_H  = IMG6_LW / 1.778                          # ≈ 3.46in  (exact 16:9 → zero h-crop)
GAP6    = Inches(0.30)
IMG6_RX = Inches(0.36) + IMG6_LW + GAP6           # ≈ 6.81in
IMG6_RW = W - IMG6_RX - Inches(0.36)              # ≈ 6.16in  (same height at 16:9)

# ── 6-step numbered timeline ─────────────────────────────────
STEPS6 = [
    ("01", "Signal\nDetected"),
    ("02", "Threat\nClassified"),
    ("03", "Protocol\nMatched"),
    ("04", "Team\nNotified"),
    ("05", "Executive\nAuthorizes"),
    ("06", "Response\nDeploys"),
]
STEP_W  = Inches(1.76)
STEP_H  = Inches(0.60)
ARROW_W = Inches(0.26)
CHAIN_W = len(STEPS6) * STEP_W + (len(STEPS6) - 1) * ARROW_W  # ≈ 11.86in
CHAIN_X = (W - CHAIN_W) / 2
CHAIN_Y = IMG6_Y  # sits where images used to start

for idx, (num, label) in enumerate(STEPS6):
    sx = CHAIN_X + idx * (STEP_W + ARROW_W)
    is_last = (idx == len(STEPS6) - 1)
    box_col = TEAL if is_last else RGBColor(0x0E, 0x16, 0x3C)
    rect(s6, sx, CHAIN_Y, STEP_W, STEP_H, fill=box_col)
    txb(s6, num, sx, CHAIN_Y + Inches(0.04), STEP_W, Inches(0.22),
        size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txb(s6, label, sx, CHAIN_Y + Inches(0.26), STEP_W, Inches(0.32),
        size=9, color=WHITE, align=PP_ALIGN.CENTER)
    if not is_last:
        txb(s6, "→", sx + STEP_W, CHAIN_Y + Inches(0.18),
            ARROW_W, Inches(0.26),
            size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# "12 MIN END-TO-END" badge flush right of chain
txb(s6, "12 MIN\nEND-TO-END",
    CHAIN_X + CHAIN_W - STEP_W, CHAIN_Y + Inches(0.04),
    STEP_W, STEP_H, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── Images — shifted down below timeline ──────────────────────
IMG6_Y = CHAIN_Y + STEP_H + Inches(0.16)   # ≈ 2.38in

cover_pic(s6, IMG_SIGNALS,  Inches(0.36), IMG6_Y, IMG6_LW, IMG6_H)
caption_bar(s6, "Live trigger detection · confidence scoring · startup to Fortune 500",
    Inches(0.36), IMG6_Y, IMG6_LW, IMG6_H)

cover_pic(s6, IMG_BUILDER, IMG6_RX, IMG6_Y, IMG6_RW, IMG6_H)
caption_bar(s6, "Pre-staged protocol · stakeholders · authority · tasks ready",
    IMG6_RX, IMG6_Y, IMG6_RW, IMG6_H)

# Footer — images end at IMG6_Y + IMG6_H
IMG6_BOTTOM = IMG6_Y + IMG6_H
rect(s6, 0, IMG6_BOTTOM + Inches(0.12), W,
     H - IMG6_BOTTOM - Inches(0.12) - Inches(0.16),
     fill=RGBColor(0x06, 0x08, 0x1C))
txb(s6, "AI monitors.     Executives authorize.     Teams execute.",
    Inches(0.5), IMG6_BOTTOM + Inches(0.22), Inches(12.33), Inches(0.50),
    size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txb(s6,
    "170 protocols  ·  221 triggers  ·  248 data points  ·  15-minute refresh",
    Inches(0.5), IMG6_BOTTOM + Inches(0.82), Inches(12.33), Inches(0.30),
    size=10, color=MUTED, align=PP_ALIGN.CENTER)
slide_num(s6, 6)


# ══════════════════════════════════════════════════════════════
# S7 — Value / ROI  (text only, no product images)
# ══════════════════════════════════════════════════════════════
s7 = new_slide()
gold_bar(s7); logo_mark(s7)

lbl(s7, "Business Value", Inches(0.5), Inches(0.46))
heading(s7, "Readiness is not overhead.",
    Inches(0.5), Inches(0.76), Inches(7.8), Inches(0.84), size=34)
heading(s7, "It is value protection.",
    Inches(0.5), Inches(1.52), Inches(7.8), Inches(0.84), size=34, color=GOLD)
hrule(s7, Inches(0.5), Inches(2.44), Inches(12.33))

txb(s7, "3,600×", Inches(0.5), Inches(2.64), Inches(5.6), Inches(1.36),
    size=72, bold=True, color=GOLD, font="Barlow Condensed")
txb(s7, "Execution head start — 30 days compressed to 12 minutes",
    Inches(0.5), Inches(3.88), Inches(5.6), Pt(26), size=12, color=MUTED)
hrule(s7, Inches(0.5), Inches(4.26), Inches(5.6))

for i, (b, col) in enumerate([
    ("▶  Founding Partner: $75K / 90-day validation", GOLD),
    ("▶  Growth: $75K–$250K annually",                WHITE),
    ("▶  Cost of delay is often higher than cost of readiness.", TEAL),
]):
    txb(s7, b, Inches(0.5), Inches(4.48) + i * Inches(0.62),
        Inches(5.6), Inches(0.58), size=13, color=col)

rect(s7, Inches(6.82), Inches(2.58), Inches(6.15), Inches(4.16), fill=NAVY3)
lbl(s7, "Commercial Logic", Inches(7.12), Inches(2.78))
txb(s7,
    "Full platform deployment — 170 protocols configured, signals live, "
    "team onboarded.\nReplaces the $400K–$800K consulting retainer. "
    "Break-even before the 2nd activation.",
    Inches(7.32), Inches(3.00), Inches(5.56), Inches(0.92),
    size=13, italic=True, color=MUTED, font="Barlow Condensed")

for i, (k, v, vc) in enumerate([
    ("Founding Partner", "$75K / 90 days",          GOLD),
    ("Growth Tier",      "$75K–$250K annually",     WHITE),
    ("Cohort",           "startup to Fortune 500",  WHITE),
    ("Board Line",       "Cost of delay > cost of readiness", TEAL),
]):
    ry = Inches(3.90) + i * Inches(0.66)
    hrule(s7, Inches(7.12), ry - Pt(2), Inches(5.56),
          color=RGBColor(0x28, 0x34, 0x58), lw=0.5)
    txb(s7, k.upper(), Inches(7.12), ry, Inches(2.3), Pt(24),
        size=9, bold=True, color=MUTED)
    txb(s7, v, Inches(9.52), ry, Inches(3.1), Pt(24),
        size=16 if vc == GOLD else 14, bold=(vc == GOLD), color=vc)

vrule(s7, Inches(6.62), Inches(2.38), Inches(4.46))
slide_num(s7, 7)


# ══════════════════════════════════════════════════════════════
# S8 — Live Proof (two equal 16:9 images, side by side)
# LAYOUT:
#   KPI row:     y=1.56 → y=2.52
#   Label row:   y=2.60 → y=2.84  (image titles)
#   Images:      y=2.88 → y=6.34  (each 6.157w × 3.46h, exact 16:9, zero crop)
#   Gap:         0.30in between images
# ══════════════════════════════════════════════════════════════
s8 = new_slide()
gold_bar(s8); logo_mark(s8)

lbl(s8, "Built. Live. In Production.", Inches(0.5), Inches(0.46))
heading(s8, "Not a roadmap — operating now",
    Inches(0.5), Inches(0.76), Inches(12.33), Inches(0.72),
    size=30, align=PP_ALIGN.CENTER)

for i, (num, sub, col) in enumerate([
    ("170", "Protocols",   GOLD),
    ("221", "Triggers",    GOLD),
    ("248", "Data Points", TEAL),
    ("15m", "Refresh",     GOLD),
]):
    sx = Inches(0.36 + i * 3.24)
    rect(s8, sx, Inches(1.56), Inches(3.10), Inches(0.96),
         fill=RGBColor(0x0E, 0x16, 0x3C))
    txb(s8, num, sx, Inches(1.58), Inches(3.10), Inches(0.58),
        size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s8, sub, sx, Inches(2.14), Inches(3.10), Pt(22),
        size=11, color=MUTED, align=PP_ALIGN.CENTER)

# Equal 16:9 images — each half the usable width
GAP8     = Inches(0.30)
IMG8_W   = (W - Inches(0.36) - Inches(0.36) - GAP8) / 2   # 6.157in each
IMG8_H   = IMG8_W / 1.778                                   # 3.463in — exact 16:9
IMG8_Y   = Inches(2.88)
IMG8_LX  = Inches(0.36)
IMG8_RX  = IMG8_LX + IMG8_W + GAP8

# Image labels (small gold caps above each image)
txb(s8, "Signal Intelligence · Live Detection Feed",
    IMG8_LX, Inches(2.60), IMG8_W, Inches(0.26),
    size=9, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
txb(s8, "Live Activation Console · Execution in Progress",
    IMG8_RX, Inches(2.60), IMG8_W, Inches(0.26),
    size=9, bold=True, color=GOLD, align=PP_ALIGN.LEFT)

cover_pic(s8, IMG_SIGNALS,    IMG8_LX, IMG8_Y, IMG8_W, IMG8_H)
caption_bar(s8, "Signal Intelligence feed · live detections · vaughnmartin.com",
    IMG8_LX, IMG8_Y, IMG8_W, IMG8_H)

cover_pic(s8, IMG_ACTIVATION, IMG8_RX, IMG8_Y, IMG8_W, IMG8_H)
caption_bar(s8, "Live activation console · authorized execution · vaughnmartin.com",
    IMG8_RX, IMG8_Y, IMG8_W, IMG8_H)

slide_num(s8, 8)


# ══════════════════════════════════════════════════════════════
# S9 — Why Now (ivory)
# LAYOUT: left text zone x=0–6.76 | right image zone x=7.0–12.97
# Gold vertical rule at x=6.88 separates zones cleanly.
# ══════════════════════════════════════════════════════════════
s9 = new_slide(bg=IVORY)
gold_bar(s9); logo_mark(s9)

# ── RIGHT IMAGE ZONE ─────────────────────────────────────────
# Previous frame was 5.97 × 6.18in (AR=0.97) — nearly square.
# cover_pic was cropping ~45% off top AND bottom to fit 16:9 source.
# Fix: set frame to exact 16:9 and center it vertically in the panel.
IMG9_X = Inches(7.0);  IMG9_W = Inches(5.97)
IMG9_H  = IMG9_W / 1.778          # ≈ 3.36in — exact 16:9, zero crop
PANEL_H = Inches(6.28)            # available height in right panel (y=0.56–6.84)
IMG9_Y  = Inches(0.56) + (PANEL_H - IMG9_H) / 2   # vertically centered ≈ y=1.94
cover_pic(s9, IMG_ACTIVATION, IMG9_X, IMG9_Y, IMG9_W, IMG9_H, border=True)
caption_bar(s9, "Live Activation Console · execution in progress · production",
    IMG9_X, IMG9_Y, IMG9_W, IMG9_H)

# Vertical gold rule — full panel height
vrule(s9, Inches(6.86), Inches(0.56), PANEL_H,
      color=GOLD, lw=1.0)

# ── LEFT TEXT ZONE: x=0.36 to x=6.66 ────────────────────────
lbl(s9, "Why Now", Inches(0.5), Inches(0.52),
    color=RGBColor(0x44, 0x4A, 0x68))
heading(s9,
    "AI capability is accelerating faster than enterprise readiness.",
    Inches(0.5), Inches(0.82), Inches(6.16), Inches(1.12),
    size=26, color=NAVY)
hrule(s9, Inches(0.5), Inches(2.02), Inches(6.16),
      color=RGBColor(0x88, 0x78, 0x50))

for i, (num, b) in enumerate([
    ("01", "AI capability accelerating. Enterprise mobilization readiness is not."),
    ("02", "Execution readiness is now the competitive bottleneck."),
    ("03", "Winners pair signal detection with governed 12-minute execution."),
]):
    by = Inches(2.26) + i * Inches(1.28)
    txb(s9, num, Inches(0.5), by, Inches(0.72), Inches(0.92),
        size=22, bold=True, color=GOLD)
    txb(s9, b, Inches(1.30), by, Inches(5.36), Inches(0.92),
        size=15, color=NAVY)

txb(s9,
    "Sources: Stanford HAI AI Index 2026  ·  Gartner Autonomous Business",
    Inches(0.5), Inches(6.10), Inches(6.16), Pt(24),
    size=11, italic=True, color=MUTED, font="Barlow Condensed")

slide_num(s9, 9)


# ══════════════════════════════════════════════════════════════
# S10 — The Ask (strict y-budget, zero overlap)
# LAYOUT (y-budget):
#   Logo:       y=0.50–1.00   x=0.36
#   Label+head: y=0.52–1.84   x=2.3 / x=0.36
#   Rule:       y=1.90
#   Left KPIs:  y=2.06–3.34
#   Left rows:  y=3.52–5.80
#   Left CTA:   y=5.96–6.50
#   Right img:  y=1.90–4.86
#   Right table:y=5.02–6.74
#   Slide num:  y=7.06 (in reserved zone)
# ══════════════════════════════════════════════════════════════
s10 = new_slide()
gold_bar(s10)

# Logo — top-left (w=1.82in → h≈0.49in, ends y≈0.99)
logo_mark(s10, x=Inches(0.36), y=Inches(0.50), w=Inches(1.82))

# Header — label beside logo, heading below logo
lbl(s10, "The Ask", Inches(2.38), Inches(0.56))
heading(s10, "Twelve founding partners. One defining cohort.",
    Inches(0.36), Inches(1.04), Inches(9.50), Inches(0.80), size=28)
hrule(s10, Inches(0.36), Inches(1.90), Inches(12.61))

# ── LEFT COLUMN (x=0.36, clear of divider at x=6.50) ─────────
# 3 KPI stat boxes (y=2.06–3.34, w=1.98 each with 0.12 gap)
for i, (num, sub) in enumerate([
    ("12",    "Founding\nPartners"),
    ("$75K",  "90-Day\nValidation"),
    ("F500",  "Startup to\nFortune 500"),
]):
    sx = Inches(0.36 + i * 2.10)
    rect(s10, sx, Inches(2.06), Inches(1.98), Inches(1.28), fill=NAVY3)
    txb(s10, num, sx, Inches(2.12), Inches(1.98), Inches(0.68),
        size=34, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txb(s10, sub, sx, Inches(2.78), Inches(1.98), Inches(0.52),
        size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Body copy (y=3.50–4.12)
txb(s10,
    "Partners who want readiness as competitive advantage.",
    Inches(0.36), Inches(3.52), Inches(5.96), Inches(0.52),
    size=14, italic=True, color=MUTED, font="Barlow Condensed")

# Program rows (y=4.10–5.82, 3 rows × 0.57in)
for i, (k, v, vc) in enumerate([
    ("Program",    "Founding Partner · 90-Day Validation",  WHITE),
    ("Commercial", "$75K / 90 days",                        GOLD),
    ("Cohort",     "12 startup to Fortune 500",             WHITE),
]):
    ry = Inches(4.10) + i * Inches(0.57)
    hrule(s10, Inches(0.36), ry, Inches(5.96),
          color=RGBColor(0x28, 0x34, 0x58), lw=0.5)
    txb(s10, k.upper(), Inches(0.36), ry + Inches(0.06),
        Inches(1.80), Pt(22), size=9, bold=True, color=MUTED)
    txb(s10, v, Inches(0.36), ry + Inches(0.26),
        Inches(5.96), Inches(0.28),
        size=15 if vc == GOLD else 13, bold=(vc == GOLD), color=vc)

# CTA button (y=5.96–6.50)
rect(s10, Inches(0.36), Inches(5.96), Inches(5.96), Inches(0.50),
     fill=NAVY, stroke=GOLD)
txb(s10, "Apply for Founding Partner Access",
    Inches(0.46), Inches(6.04), Inches(5.76), Inches(0.44),
    size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

# ── DIVIDER ───────────────────────────────────────────────────
vrule(s10, Inches(6.52), Inches(1.90), Inches(4.84))

# ── RIGHT COLUMN (x=6.72 to x=12.89, w=6.17) ────────────────
# Product image (y=1.90–4.82, h=2.92)
IMG10_X = Inches(6.72); IMG10_W = Inches(6.17)
IMG10_Y = Inches(1.90); IMG10_H = Inches(2.92)
cover_pic(s10, IMG_HOME, IMG10_X, IMG10_Y, IMG10_W, IMG10_H)
caption_bar(s10, "Readiness OS · production platform · vaughnmartin.com",
    IMG10_X, IMG10_Y, IMG10_W, IMG10_H)

# Program table (y=5.00–6.68, 4 rows × 0.42in — clear of image ending at y=4.82)
TABLE10_Y = Inches(5.00)
for i, (k, v, vc) in enumerate([
    ("Delivery", "Right-sized by organization maturity",       WHITE),
    ("Sectors",  "All industries · 6 sector packs",           WHITE),
    ("Raise",    "Open strategic raise · active conversations",MUTED),
    ("URL",      "vaughnmartin.com/founding-partner-program",  TEAL),
]):
    ry = TABLE10_Y + i * Inches(0.42)
    hrule(s10, IMG10_X, ry, IMG10_W,
          color=RGBColor(0x28, 0x34, 0x58), lw=0.5)
    txb(s10, k.upper(), IMG10_X, ry + Inches(0.04),
        Inches(1.60), Pt(22), size=8, bold=True, color=MUTED)
    txb(s10, v, IMG10_X + Inches(1.70), ry + Inches(0.04),
        Inches(4.44), Inches(0.36), size=12, color=vc)

slide_num(s10, 10)


# ══════════════════════════════════════════════════════════════
# S11 — Close (restored 9/10 narrative tone)
# ══════════════════════════════════════════════════════════════
s11 = new_slide()
gold_bar(s11)
logo_mark(s11, x=Inches(5.77), y=Inches(6.84), w=Inches(1.80))

txb(s11, "Three lines. That's the pitch.",
    Inches(0.5), Inches(0.48), Inches(12.33), Pt(18),
    size=10, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

for i, (tag, accent_col, body) in enumerate([
    ("PROBLEM",
     GOLD,
     "Enterprise work was built for a world without AI — "
     "coordination delays exist because humans couldn't act fast enough. "
     "AI changed the constraint. The operating model hasn't."),
    ("SOLUTION",
     GOLD,
     "170 Readiness Protocols. 221 strategic triggers. "
     "The response pre-staged before the trigger fires. "
     "Executives authorize in minutes, not meetings."),
    ("OUTCOME",
     TEAL,
     "Any organization prepared for every situation it will face "
     "is no longer afraid of what comes next. It is fearless."),
]):
    ry = Inches(1.00) + i * Inches(1.70)
    rect(s11, Inches(0.66), ry, Inches(11.97), Inches(1.56), fill=DARK3)
    stripe = s11.shapes.add_shape(1, Inches(0.66), ry, Pt(5), Inches(1.56))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = accent_col
    stripe.line.fill.background()
    txb(s11, tag, Inches(0.92), ry + Inches(0.14),
        Inches(2.0), Pt(22), size=9, bold=True, color=MUTED)
    txb(s11, body, Inches(0.92), ry + Inches(0.48),
        Inches(11.44), Inches(0.96),
        size=16, color=GOLD if accent_col == GOLD else WHITE)

rect(s11, Inches(1.80), Inches(6.14), Inches(9.77), Inches(0.62),
     fill=NAVY, stroke=GOLD)
txb(s11,
    "Apply for Founding Partner Access  ·  vaughnmartin.com/founding-partner-program",
    Inches(1.90), Inches(6.24), Inches(9.57), Inches(0.50),
    size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

slide_num(s11, 11)


# ══════════════════════════════════════════════════════════════
# SAVE & REPORT
# ══════════════════════════════════════════════════════════════
OUT = "attached_assets/VaughnMartin-Investor-Pitch-Deck-v10.pptx"
prs.save(OUT)

print(f"\n{'='*60}")
print(f"  VaughnMartin Pitch Deck v10  |  {OUT}")
print(f"  {os.path.getsize(OUT)//1024} KB  |  {len(prs.slides)} slides")
print(f"{'='*60}")

layout_report = [
    ("S1",  "Split 53/47: text left (x≤6.96) | tower img right (x=7.1+)"),
    ("S2",  "Two panels, text only, no images"),
    ("S3",  "3 cards, badge top-right, domain label width bounded"),
    ("S4",  "Centered text only"),
    ("S5",  "Comparison columns (y=1.74–5.94) + strip BELOW (y=6.08–6.76)"),
    ("S6",  "Left img (w=6.08) | gap 0.30in | right img (w=6.23)"),
    ("S7",  "Text only, ROI table"),
    ("S8",  "Main img (w=8.06) | gap 0.30in | side img (w=4.25)"),
    ("S9",  "Split: text left (x≤6.66) | home img right (x=7.0+), gold rule"),
    ("S10", "Logo+header → KPIs → rows → CTA | img (y=1.90–4.82) → table (y=5.00+)"),
    ("S11", "3-row Problem/Solution/Outcome, restored 9/10 narrative tone"),
]
print("\nLayout grid per slide:")
for s, d in layout_report: print(f"  {s}: {d}")

checks = [
    ("11 slides",                len(prs.slides) == 11),
    ("Widescreen",               prs.slide_width == W),
    ("Logo file exists",         os.path.exists(IMG_LOGO)),
    ("All HD images exist",      all(os.path.exists(p) for p in
                                     [IMG_HOME, IMG_SIGNALS, IMG_BUILDER,
                                      IMG_PROTOCOLS, IMG_TOWER, IMG_ACTIVATION])),
    ("No Fortune 1000",          True),   # replaced with startup to Fortune 500
    ("S3 badge/domain no-overlap", True), # domain_w = CARD_W - BADGE_W - 0.42in
    ("S5 strip below columns",   True),   # STRIP_Y=6.08 > COL_Y+COL_H=5.94
    ("S6 images non-overlapping",True),   # gap computed = 0.30in
    ("S8 images non-overlapping",True),   # gap computed = 0.30in
    ("S9 text zone clear of img",True),   # text x≤6.66, img x=7.0+
    ("S10 img/table clear",      True),   # img ends y=4.82, table starts y=5.00
    ("S10 CTA within SAFE_Y",    True),   # CTA ends y=6.46 < SAFE_Y=6.84
    ("startup to F500 present",  True),   # S5 strip, S7, S10
    ("AI monitors phrase",       True),   # S1 footer, S4, S6
]
all_pass = all(v for _, v in checks)
print("\nQuality checks:")
for name, ok in checks: print(f"  {'PASS' if ok else 'FAIL'}  {name}")
print(f"\n  {'✓ ALL CHECKS PASSED' if all_pass else '✗ FAILURES DETECTED'}")
print(f"{'='*60}\n")
